"""Build realistic war-room fixtures (md / pdf / pptx / png)."""

from __future__ import annotations

import struct
import zlib
from pathlib import Path

from _common import FIXTURES, OUTPUT


def _minimal_png(width: int = 320, height: int = 180, bar_height: int = 90) -> bytes:
    """Synthetic 'dashboard spike' PNG (red error bar on dark background)."""

    def chunk(tag: bytes, data: bytes) -> bytes:
        return (
            struct.pack(">I", len(data))
            + tag
            + data
            + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF)
        )

    rows = []
    for y in range(height):
        row = [0]  # filter
        for x in range(width):
            # dark bg
            r, g, b = 18, 22, 28
            # axes
            if y == height - 20 or x == 30:
                r = g = b = 90
            # error spike bar
            if 140 <= x <= 200 and (height - 20 - bar_height) <= y < height - 20:
                r, g, b = 220, 50, 50
            # green baseline
            elif 40 <= x < 140 and (height - 35) <= y < height - 20:
                r, g, b = 40, 170, 90
            row.extend([r, g, b])
        rows.append(bytes(row))
    raw = b"".join(rows)
    ihdr = struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0)
    return (
        b"\x89PNG\r\n\x1a\n"
        + chunk(b"IHDR", ihdr)
        + chunk(b"IDAT", zlib.compress(raw, 9))
        + chunk(b"IEND", b"")
    )


def _write_pdf(path: Path, lines: list[str]) -> None:
    content = ["BT", "/F1 11 Tf", "50 750 Td"]
    for i, line in enumerate(lines):
        safe = line.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
        if i:
            content.append("0 -14 Td")
        content.append(f"({safe}) Tj")
    content.append("ET")
    stream = "\n".join(content)
    objects = [
        "1 0 obj<< /Type /Catalog /Pages 2 0 R >>endobj\n",
        "2 0 obj<< /Type /Pages /Kids [3 0 R] /Count 1 >>endobj\n",
        (
            "3 0 obj<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            "/Contents 4 0 R /Resources<< /Font<< /F1 5 0 R >> >> >>endobj\n"
        ),
        f"4 0 obj<< /Length {len(stream)} >>stream\n{stream}\nendstream\nendobj\n",
        "5 0 obj<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>endobj\n",
    ]
    header = "%PDF-1.4\n"
    offsets = [0]
    pos = len(header)
    body = ""
    for obj in objects:
        offsets.append(pos)
        body += obj
        pos += len(obj.encode("latin-1"))
    xref = f"xref\n0 {len(offsets)}\n0000000000 65535 f \n"
    for off in offsets[1:]:
        xref += f"{off:010d} 00000 n \n"
    trailer = (
        f"trailer<< /Size {len(offsets)} /Root 1 0 R >>\nstartxref\n{pos}\n%%EOF\n"
    )
    path.write_bytes((header + body + xref + trailer).encode("latin-1"))


def _write_pptx(path: Path) -> None:
    from pptx import Presentation

    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "AcmePay Incident Status"
    s1.placeholders[1].text = "INC-88421 — Settlement Rail v3 — SEV candidate"

    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "Timeline (IST)"
    s2.placeholders[1].text = (
        "18:40 — elevated settlement error rate in ap-south-1\n"
        "18:42 — BharatNova NOC pages AcmePay\n"
        "18:51 — connector pool saturation suspected\n"
        "19:05 — war room opened; customer bridge pending"
    )

    s3 = prs.slides.add_slide(prs.slide_layouts[1])
    s3.shapes.title.text = "Impact & Next Actions"
    s3.placeholders[1].text = (
        "Impact: ~42k merchant payouts delayed\n"
        "Hypothesis: bank connector thread pool exhaustion after cert rotation\n"
        "Next: freeze deploys, scale connector workers, dual-write audit log"
    )
    prs.save(str(path))


def main() -> None:
    FIXTURES.mkdir(parents=True, exist_ok=True)
    OUTPUT.mkdir(parents=True, exist_ok=True)

    (FIXTURES / "runbook.md").write_text(
        """# Settlement Rail v3 — SEV Runbook

## Symptoms
- Partner payouts stuck in RETRYING/FAILED
- `settlement.batch.submit` error spike
- Connector pool wait time > 2s

## Immediate actions
1. Confirm region health (`ap-south-1` settlement + connector)
2. Check latest cert/config change in last 6h
3. If connector saturation: scale workers, drain queue, notify partner
4. Strategic tier: update within **15 minutes**, bridge within **30 minutes**

## Communication
- Use customer-safe language (no internal hostnames)
- State ETA + mitigation clearly
""",
        encoding="utf-8",
    )

    _write_pdf(
        FIXTURES / "strategic_sla.pdf",
        [
            "AcmePay - Strategic Partner SLA (excerpt)",
            "Partner tier: Strategic",
            "P1 response target: 15 minutes",
            "P1 restore target: 60 minutes",
            "Customer bridge: required within 30 minutes of P1 ack",
            "Credits: 5% monthly fee per qualifying P1 breach hour",
            "Region coverage: ap-south-1 production settlement rail",
        ],
    )
    _write_pptx(FIXTURES / "incident_status.pptx")
    (FIXTURES / "dashboard_spike.png").write_bytes(_minimal_png())
    print(f"Fixtures ready under {FIXTURES}")


if __name__ == "__main__":
    main()
