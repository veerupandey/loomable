"""04 — Document I/O: markdown, PDF, PowerPoint

USE WHEN: The agent should read documents (md/pdf/pptx) and write outputs.

Creates sample files, then asks the agent to:
  1) read markdown + PDF + PPT
  2) write a combined markdown summary
  3) (optional) write a short brief back to disk
"""

from __future__ import annotations

import asyncio
import sys
import tempfile
from pathlib import Path

_HERE = Path(__file__).resolve().parent
if str(_HERE) not in sys.path:
    sys.path.insert(0, str(_HERE))

from loomable.agent import Agent
from loomable.toolkits import FileTools, PDFTools, PPTTools

from _provider import make_provider


def _write_minimal_pdf(path: Path, lines: list[str]) -> None:
    """Write a tiny text PDF without extra deps (Helvetica)."""
    # Escape parentheses for PDF text operators.
    content_lines = ["BT", "/F1 12 Tf", "50 750 Td"]
    for i, line in enumerate(lines):
        safe = line.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
        if i == 0:
            content_lines.append(f"({safe}) Tj")
        else:
            content_lines.append("0 -16 Td")
            content_lines.append(f"({safe}) Tj")
    content_lines.append("ET")
    stream = "\n".join(content_lines)
    objects = []
    objects.append("1 0 obj<< /Type /Catalog /Pages 2 0 R >>endobj\n")
    objects.append("2 0 obj<< /Type /Pages /Kids [3 0 R] /Count 1 >>endobj\n")
    objects.append(
        "3 0 obj<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        "/Contents 4 0 R /Resources<< /Font<< /F1 5 0 R >> >> >>endobj\n"
    )
    objects.append(
        f"4 0 obj<< /Length {len(stream)} >>stream\n{stream}\nendstream\nendobj\n"
    )
    objects.append(
        "5 0 obj<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>endobj\n"
    )

    body = "".join(objects)
    # Build xref
    header = "%PDF-1.4\n"
    offsets = [0]
    offset = len(header)
    for obj in objects:
        offsets.append(offset)
        offset += len(obj.encode("latin-1"))
    xref_pos = offset
    xref = [f"xref\n0 {len(offsets)}\n", "0000000000 65535 f \n"]
    for off in offsets[1:]:
        xref.append(f"{off:010d} 00000 n \n")
    trailer = (
        f"trailer<< /Size {len(offsets)} /Root 1 0 R >>\n"
        f"startxref\n{xref_pos}\n%%EOF\n"
    )
    path.write_bytes(
        (header + body + "".join(xref) + trailer).encode("latin-1")
    )


def _write_sample_pptx(path: Path) -> None:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Q2 Ops Update"
    slide.placeholders[1].text = (
        "India expansion on track.\n"
        "Hiring 12 engineers in Bengaluru.\n"
        "Risk: delayed compliance review."
    )
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Next Actions"
    slide2.placeholders[1].text = (
        "1. Finish SOC2 evidence pack\n"
        "2. Publish customer case study\n"
        "3. Lock pricing for mid-market"
    )
    prs.save(str(path))


def _build_workspace() -> Path:
    root = Path(tempfile.mkdtemp(prefix="loomable_docs_"))
    (root / "notes.md").write_text(
        "# Internal Notes\n\n"
        "- Product: AI shop-floor scheduling\n"
        "- Target: mid-size factories in India\n"
        "- Ask: summarize all docs into one brief\n",
        encoding="utf-8",
    )
    _write_minimal_pdf(
        root / "market.pdf",
        [
            "Market Snapshot",
            "India manufacturing digitization is rising.",
            "Buyers care about uptime and simple pricing.",
            "Competitors: legacy MES and spreadsheets.",
        ],
    )
    _write_sample_pptx(root / "ops.pptx")
    return root


async def main() -> None:
    workspace = _build_workspace()
    print(f"Workspace: {workspace}")
    print("Files:", ", ".join(p.name for p in workspace.iterdir()))

    agent = Agent(
        model=make_provider(),
        role="Document analyst",
        goal="Read md/pdf/pptx and write a clear combined brief",
        instructions=(
            f"Your working directory base is configured. "
            f"Read notes.md, market.pdf, and ops.pptx. "
            f"Then write output/summary.md with a short combined brief. "
            f"Be concrete and plain English."
        ),
        tools=[
            FileTools(base_dir=str(workspace)),
            PDFTools(),
            PPTTools(),
        ],
    )

    # PDFTools/PPTTools use absolute paths; FileTools uses base_dir-relative.
    result = await agent.arun(
        "Please read these inputs and produce one markdown summary:\n"
        f"1) markdown file: notes.md\n"
        f"2) pdf file: {workspace / 'market.pdf'}\n"
        f"3) pptx file: {workspace / 'ops.pptx'}\n"
        "Write the final brief to output/summary.md using write_file."
    )
    print("\n=== Agent answer ===\n")
    print(result.output.text())
    summary = workspace / "output" / "summary.md"
    if summary.exists():
        print("\n=== output/summary.md ===\n")
        print(summary.read_text(encoding="utf-8"))
    else:
        print("\n[note] summary.md was not written — check tool activity")
    if result.tool_activity:
        print(f"\n[tools used: {len(result.tool_activity)}]")


if __name__ == "__main__":
    asyncio.run(main())
