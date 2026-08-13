"""loomable.toolkits.ppt_tools - PowerPoint reading toolkit.

Provides tools for extracting text from ``.pptx`` files.
Requires the ``python-pptx`` package (install via ``pip install loomable[ppt]``
or ``pip install python-pptx``).
"""

from __future__ import annotations

import asyncio
from pathlib import Path

from loomable.agent.tools import FunctionTool
from loomable.toolkits._base import Toolkit


class PPTTools(Toolkit):
    """PowerPoint (.pptx) reading toolkit. Requires: python-pptx"""

    def __init__(
        self,
        *,
        include_tools: list[str] | None = None,
        exclude_tools: list[str] | None = None,
    ) -> None:
        try:
            import pptx  # noqa: F401
        except ImportError as exc:
            raise ImportError(
                "PPTTools requires 'python-pptx'. "
                "Install with: pip install python-pptx"
            ) from exc
        super().__init__(include_tools=include_tools, exclude_tools=exclude_tools)

    def _register_tools(self) -> list[FunctionTool]:
        return [
            FunctionTool(self._read_pptx, name="read_pptx"),
            FunctionTool(self._list_pptx_slides, name="list_pptx_slides"),
        ]

    async def _read_pptx(self, path: str, slides: str | None = None) -> str:
        """Read text from a PowerPoint file. Optional slides e.g. '1-3' or '1,3'."""
        return await asyncio.to_thread(self._read_pptx_sync, path, slides)

    async def _list_pptx_slides(self, path: str) -> str:
        """List slide numbers and titles/first lines from a PowerPoint file."""
        return await asyncio.to_thread(self._list_pptx_slides_sync, path)

    def _read_pptx_sync(self, path: str, slides: str | None) -> str:
        from pptx import Presentation

        file_path = Path(path)
        if not file_path.exists():
            return f"Error: File not found: {path}"

        prs = Presentation(str(file_path))
        indices = self._parse_slides(slides, len(prs.slides))
        chunks: list[str] = []
        for i in indices:
            slide = prs.slides[i]
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            body = "\n".join(texts) if texts else "(no text)"
            chunks.append(f"--- Slide {i + 1} ---\n{body}")
        return "\n\n".join(chunks) if chunks else "(empty presentation)"

    def _list_pptx_slides_sync(self, path: str) -> str:
        from pptx import Presentation

        file_path = Path(path)
        if not file_path.exists():
            return f"Error: File not found: {path}"

        prs = Presentation(str(file_path))
        lines: list[str] = []
        for i, slide in enumerate(prs.slides, start=1):
            title = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    title = shape.text.strip().splitlines()[0][:80]
                    break
            lines.append(f"{i}. {title or '(untitled)'}")
        return "\n".join(lines) if lines else "(no slides)"

    def _parse_slides(self, slides: str | None, total: int) -> list[int]:
        if not slides:
            return list(range(total))
        slides = slides.strip()
        out: list[int] = []
        try:
            if "-" in slides and "," not in slides:
                start_s, end_s = slides.split("-", 1)
                start, end = int(start_s), int(end_s)
                out = list(range(max(1, start) - 1, min(total, end)))
            else:
                for part in slides.split(","):
                    n = int(part.strip())
                    if 1 <= n <= total:
                        out.append(n - 1)
        except ValueError:
            return list(range(total))
        return out or list(range(total))
