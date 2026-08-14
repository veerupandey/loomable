"""Load heterogeneous sources into :class:`~loomable.retrieval.types.Document`s.

Default ingest covers popular docs, code, HTML, PDF, Office, JSON/CSV, and
``http(s)`` URLs. Optional deps: ``pypdf`` (PDF), ``beautifulsoup4`` (HTML),
``python-pptx`` (PPTX). DOCX uses the stdlib zip/XML reader.
"""

from __future__ import annotations

import json
import mimetypes
import re
import zipfile
from pathlib import Path
from typing import Any, Iterable, Sequence
from urllib.parse import urlparse
from xml.etree import ElementTree as ET

from loomable.retrieval.types import Document

_SKIP_DIRS = {
    ".git",
    ".hg",
    ".svn",
    ".venv",
    "venv",
    "node_modules",
    "__pycache__",
    ".mypy_cache",
    ".pytest_cache",
    ".ruff_cache",
    ".tox",
    "dist",
    "build",
    ".deep_workspace",
    ".sandbox",
    ".loomable",
    ".next",
    "target",
    "vendor",
}

_TEXT_EXTS = {
    ".txt",
    ".md",
    ".mdx",
    ".rst",
    ".adoc",
    ".py",
    ".pyi",
    ".ts",
    ".tsx",
    ".js",
    ".jsx",
    ".mjs",
    ".cjs",
    ".go",
    ".rs",
    ".java",
    ".kt",
    ".kts",
    ".cs",
    ".rb",
    ".php",
    ".swift",
    ".scala",
    ".sc",
    ".c",
    ".h",
    ".cpp",
    ".hpp",
    ".cc",
    ".cxx",
    ".m",
    ".mm",
    ".html",
    ".htm",
    ".css",
    ".scss",
    ".less",
    ".json",
    ".jsonl",
    ".toml",
    ".yaml",
    ".yml",
    ".xml",
    ".sql",
    ".sh",
    ".bash",
    ".zsh",
    ".ps1",
    ".csv",
    ".tsv",
    ".ipynb",
    ".r",
    ".jl",
    ".lua",
    ".pl",
    ".pm",
    ".ex",
    ".exs",
    ".erl",
    ".hs",
    ".dart",
    ".vue",
    ".svelte",
    ".tf",
    ".proto",
    ".graphql",
    ".gql",
}

_LARGE_OK_EXTS = {".pdf", ".docx", ".pptx"}
_URL_RE = re.compile(r"^https?://", re.I)
_W_NS = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"


def _guess_media_type(path: Path) -> str:
    mt, _ = mimetypes.guess_type(str(path))
    return mt or "text/plain"


def is_http_url(value: str) -> bool:
    """True when ``value`` looks like an http(s) URL (not a local path)."""
    s = (value or "").strip()
    if not _URL_RE.match(s):
        return False
    parsed = urlparse(s)
    return bool(parsed.scheme in {"http", "https"} and parsed.netloc)


def _read_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader  # type: ignore
    except Exception as exc:  # noqa: BLE001
        raise RuntimeError(
            "PDF ingest requires optional dependency 'pypdf' (pip install loomable[pdf])"
        ) from exc
    reader = PdfReader(str(path))
    parts: list[str] = []
    for i, page in enumerate(reader.pages):
        try:
            text = page.extract_text() or ""
        except Exception:  # noqa: BLE001
            text = ""
        parts.append(f"--- Page {i + 1} ---\n{text}")
    return "\n".join(parts)


def _read_docx(path: Path) -> str:
    """Extract plain text from a .docx via stdlib zip/XML (no python-docx needed)."""
    try:
        with zipfile.ZipFile(path) as zf:
            xml = zf.read("word/document.xml")
    except (KeyError, zipfile.BadZipFile) as exc:
        raise RuntimeError(f"invalid docx: {path}") from exc
    root = ET.fromstring(xml)
    paras: list[str] = []
    for p in root.iter(f"{_W_NS}p"):
        texts = [(t.text or "") for t in p.iter(f"{_W_NS}t")]
        line = "".join(texts).strip()
        if line:
            paras.append(line)
    return "\n\n".join(paras)


def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as exc:  # noqa: BLE001
        raise RuntimeError(
            "PPTX ingest requires optional dependency 'python-pptx' "
            "(pip install loomable[ppt])"
        ) from exc
    prs = Presentation(str(path))
    slides: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        bits: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text and str(text).strip():
                bits.append(str(text).strip())
        if bits:
            slides.append(f"--- Slide {i} ---\n" + "\n".join(bits))
    return "\n\n".join(slides)


def _read_ipynb(path: Path) -> str:
    raw = path.read_text(encoding="utf-8", errors="replace")
    try:
        nb = json.loads(raw)
    except json.JSONDecodeError:
        return raw
    parts: list[str] = []
    for i, cell in enumerate(nb.get("cells") or []):
        src = cell.get("source") or ""
        if isinstance(src, list):
            src = "".join(src)
        cell_type = cell.get("cell_type") or "code"
        parts.append(f"--- Cell {i} ({cell_type}) ---\n{src}")
    return "\n\n".join(parts) if parts else raw


def load_url(
    url: str,
    *,
    doc_id: str | None = None,
    timeout: float = 30.0,
    block_private_hosts: bool = True,
) -> Document:
    """Fetch an http(s) URL and return a Document (HTML stripped when possible)."""
    from loomable.toolkits.net_safety import validate_http_url

    raw = (url or "").strip()
    err = validate_http_url(raw, block_private_hosts=block_private_hosts)
    if err:
        raise ValueError(err)

    import httpx

    with httpx.Client(
        timeout=timeout,
        follow_redirects=True,
        headers={"User-Agent": "loomable-ingest/0.2"},
    ) as client:
        resp = client.get(raw)
        resp.raise_for_status()
        content_type = (resp.headers.get("content-type") or "text/html").split(";")[0].strip()
        body = resp.text

    text = body
    media = content_type or "text/html"
    if "html" in media or body.lstrip().lower().startswith(("<!doctype", "<html")):
        try:
            from bs4 import BeautifulSoup  # type: ignore

            soup = BeautifulSoup(body, "html.parser")
            for tag in soup(["script", "style", "noscript"]):
                tag.decompose()
            text = soup.get_text("\n")
            media = "text/html"
        except Exception:  # noqa: BLE001
            text = re.sub(r"<[^>]+>", " ", body)
            media = "text/html"
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    return Document(
        id=doc_id or raw,
        text=text,
        source=raw,
        media_type=media,
        metadata={"url": raw, "content_type": content_type},
    )


def load_file(path: str | Path, *, doc_id: str | None = None) -> Document:
    """Load a single file into a Document (docs/code/html/pdf/docx/pptx/json/…)."""
    p = Path(path).expanduser().resolve()
    if not p.is_file():
        raise FileNotFoundError(str(p))
    media = _guess_media_type(p)
    suffix = p.suffix.lower()
    if suffix == ".pdf":
        text = _read_pdf(p)
        media = "application/pdf"
    elif suffix == ".docx":
        text = _read_docx(p)
        media = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    elif suffix == ".pptx":
        text = _read_pptx(p)
        media = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    elif suffix == ".ipynb":
        text = _read_ipynb(p)
        media = "application/x-ipynb+json"
    else:
        text = p.read_text(encoding="utf-8", errors="replace")
    return Document(
        id=doc_id or p.as_posix(),
        text=text,
        source=p.as_posix(),
        media_type=media,
        metadata={"filename": p.name, "suffix": suffix},
    )


def load_directory(
    root: str | Path,
    *,
    extensions: Sequence[str] | None = None,
    max_files: int = 5_000,
) -> list[Document]:
    """Load text-like / office / PDF files under ``root`` (skips venv/git/build)."""
    base = Path(root).expanduser().resolve()
    if not base.is_dir():
        raise NotADirectoryError(str(base))
    default_exts = (*_TEXT_EXTS, ".pdf", ".docx", ".pptx")
    exts = {
        e if e.startswith(".") else f".{e}"
        for e in (extensions or default_exts)
    }
    out: list[Document] = []
    for path in sorted(base.rglob("*")):
        if not path.is_file():
            continue
        rel_parts = path.relative_to(base).parts[:-1]
        if any(part in _SKIP_DIRS or part.startswith(".") for part in rel_parts):
            continue
        suffix = path.suffix.lower()
        if suffix not in exts and path.name.lower() != "dockerfile":
            continue
        try:
            size = path.stat().st_size
            if size > 2_000_000 and suffix not in _LARGE_OK_EXTS:
                continue
        except OSError:
            continue
        try:
            out.append(load_file(path))
        except Exception:  # noqa: BLE001 — skip unreadable
            continue
        if len(out) >= max_files:
            break
    return out


def coerce_source(source: Any) -> list[Document]:
    """Normalize a source spec into one or more Documents.

    Accepts:
    - :class:`Document`
    - ``Path`` / existing path string → file or directory
    - ``http(s)://…`` URL string → fetched page
    - ``dict`` with ``text`` or ``url`` (+ optional ``id`` / ``source`` / ``media_type``)
    - raw ``str`` that is not a path/URL → treated as inline text
    """
    if isinstance(source, Document):
        return [source]
    if isinstance(source, dict):
        extra_meta = dict(source.get("metadata") or {})
        for lift in ("title", "author", "tags", "department", "locale"):
            if source.get(lift) is not None:
                extra_meta.setdefault(lift, source[lift])
        if source.get("url"):
            docs = [
                load_url(
                    str(source["url"]),
                    doc_id=str(source["id"]) if source.get("id") else None,
                )
            ]
            for d in docs:
                d.metadata.update({k: v for k, v in extra_meta.items() if v is not None})
            return docs
        if source.get("path") or source.get("file"):
            raw_path = source.get("path") or source.get("file")
            loaded = coerce_source(raw_path)
            for d in loaded:
                d.metadata.update({k: v for k, v in extra_meta.items() if v is not None})
                if source.get("id") and len(loaded) == 1:
                    d.id = str(source["id"])
            return loaded
        text = str(source.get("text") or "")
        return [
            Document(
                id=str(source.get("id") or source.get("source") or f"doc-{hash(text) & 0xffff}"),
                text=text,
                source=str(source.get("source") or ""),
                media_type=str(source.get("media_type") or "text/plain"),
                metadata=extra_meta,
            )
        ]
    if isinstance(source, Path) or (
        isinstance(source, str) and (Path(source).expanduser().exists())
    ):
        path = Path(source).expanduser()
        if path.is_dir():
            return load_directory(path)
        if path.is_file():
            return [load_file(path)]
    if isinstance(source, str) and is_http_url(source):
        return [load_url(source)]
    if isinstance(source, str):
        return [
            Document(
                id=f"inline-{hash(source) & 0xfffffff:x}",
                text=source,
                source="inline",
                media_type="text/plain",
            )
        ]
    raise TypeError(f"unsupported retrieval source type: {type(source)!r}")


def load_sources(sources: Iterable[Any]) -> list[Document]:
    """Flatten many source specs into Documents (dedupe by id)."""
    docs: list[Document] = []
    seen: set[str] = set()
    for raw in sources:
        for doc in coerce_source(raw):
            if doc.id in seen:
                continue
            seen.add(doc.id)
            docs.append(doc)
    return docs
