"""Unit tests for PPTTools."""

from __future__ import annotations

import sys
from pathlib import Path
from unittest.mock import patch

import pytest

pytest.importorskip("pptx")

from loomable.toolkits.ppt_tools import PPTTools


def _make_pptx(path: Path, slides: list[tuple[str, str]]) -> Path:
    from pptx import Presentation

    prs = Presentation()
    for title, body in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = body
    prs.save(str(path))
    return path


@pytest.fixture
def ppt_tools() -> PPTTools:
    return PPTTools()


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
    return _make_pptx(
        tmp_path / "sample.pptx",
        [
            ("Intro", "Matcha is green tea powder."),
            ("Uses", "Drinks, desserts, wellness."),
        ],
    )


class TestImportError:
    def test_raises_import_error_when_pptx_missing(self) -> None:
        with patch.dict(sys.modules, {"pptx": None}):
            with pytest.raises(ImportError, match="python-pptx"):
                PPTTools()


class TestPPTToolsRead:
    async def test_read_all_slides(
        self, ppt_tools: PPTTools, sample_pptx: Path
    ) -> None:
        text = await ppt_tools._read_pptx(str(sample_pptx))
        assert "Slide 1" in text
        assert "Matcha" in text
        assert "Slide 2" in text
        assert "Drinks" in text

    async def test_read_slide_range(
        self, ppt_tools: PPTTools, sample_pptx: Path
    ) -> None:
        text = await ppt_tools._read_pptx(str(sample_pptx), slides="2")
        assert "Slide 2" in text
        assert "Slide 1" not in text

    async def test_list_slides(
        self, ppt_tools: PPTTools, sample_pptx: Path
    ) -> None:
        listing = await ppt_tools._list_pptx_slides(str(sample_pptx))
        assert "1. Intro" in listing
        assert "2. Uses" in listing

    async def test_missing_file(self, ppt_tools: PPTTools) -> None:
        text = await ppt_tools._read_pptx("/tmp/does-not-exist-loomable.pptx")
        assert "Error: File not found" in text

    def test_registered_tool_names(self, ppt_tools: PPTTools) -> None:
        names = {t.name for t in ppt_tools.tools()}
        assert names == {"read_pptx", "list_pptx_slides"}
